"""Assemble the talk deck into a .pptx.

Reuses the exact slide-drawing functions from plot_talk_slides (so the .pptx is
pixel-identical to the verified PDF), renders each slide to a high-res PNG, and
places one full-bleed image per 16:9 slide.

    python scripts/build_talk_pptx.py   ->  docs/talk_slides.pptx
"""
from __future__ import annotations
import sys
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))
import plot_talk_slides as T  # noqa: E402
import plot_talk_slides_compact as C  # noqa: E402

TMP = ROOT / "docs" / "_slide_png"


def render_full():
    order = [T.s1, T.s2, T.s4, T.s5, None, T.s3, T.s7]  # matches plot_talk_slides.main
    paths = []
    for i, fn in enumerate(order, start=1):
        fig, ax = T.page()
        if i == 5:
            T.s6(ax, fig)
        else:
            fn(ax)
        T.deco(ax, i)
        p = TMP / f"full_{i}.png"
        fig.savefig(p, dpi=200, facecolor=T.GROUND); plt.close(fig); paths.append(p)
    return paths


def render_compact():
    C.T.NPAGES = 3
    paths = []
    for i, fn in enumerate([C.c1, C.c2, C.c3], start=1):
        fig, ax = C.page()
        if i == 2:
            C.c2(ax, fig)
        else:
            fn(ax)
        C.T.deco(ax, i)
        p = TMP / f"compact_{i}.png"
        fig.savefig(p, dpi=200, facecolor=C.GROUND); plt.close(fig); paths.append(p)
    return paths


def assemble(pngs, out):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for p in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(out)
    print(f"saved {out}  ({len(pngs)} slides)")


def main():
    TMP.mkdir(parents=True, exist_ok=True)
    assemble(render_full(), ROOT / "docs" / "talk_slides.pptx")
    assemble(render_compact(), ROOT / "docs" / "talk_slides_compact.pptx")


if __name__ == "__main__":
    main()
